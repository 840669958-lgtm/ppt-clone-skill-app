"""
quick_test.py — ppt_analyzer 快速验证脚本
用法：
  python quick_test.py <your_file.pptx>         # 输出人类可读摘要
  python quick_test.py <your_file.pptx> --json  # 输出完整 JSON

如果手头没有 PPT，脚本会自动生成一个最简 .pptx 进行内部验证。
"""
import sys
import pathlib
import json

ROOT = pathlib.Path(__file__).parent
sys.path.insert(0, str(ROOT))

from core.ppt_analyzer import analyze, analyze_to_json, _print_summary


def _make_demo_pptx(dest: pathlib.Path) -> None:
    """生成一个带主题色、母版、多版式的演示 pptx，用于无文件时自测。"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    # 封面页
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "品牌模板演示"
    slide.placeholders[1].text = "副标题占位符"
    # 内容页
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "内容页"
    slide2.placeholders[1].text = "• 条目一\n• 条目二\n• 条目三"
    # 空白页
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(str(dest))
    print(f"[INFO] 已生成演示 PPT：{dest}")


def main() -> None:
    if len(sys.argv) >= 2 and sys.argv[1] != "--json":
        target = pathlib.Path(sys.argv[1])
    else:
        # 自动生成演示文件
        target = ROOT / "demo_template.pptx"
        if not target.exists():
            _make_demo_pptx(target)

    output_json = "--json" in sys.argv

    try:
        if output_json:
            result = analyze_to_json(target)
            print(result)
            # 写入同目录供调试
            out_file = target.with_suffix(".profile.json")
            out_file.write_text(result, encoding="utf-8")
            print(f"\n[INFO] JSON 已写入：{out_file}", file=sys.stderr)
        else:
            profile = analyze(target)
            _print_summary(profile)
    except Exception as e:
        print(f"[ERROR] {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
