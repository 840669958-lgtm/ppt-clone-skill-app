"""
core/ppt_analyzer.py
====================
PPT 模板解析引擎 —— 整个项目的数据源头。

解析维度：
  1. 页面几何：宽度 / 高度（EMU & 厘米）
  2. 母版结构：每个 SlideMaster 的背景色、装饰 Shape 数量、Logo 候选
  3. 版式清单：每个 SlideLayout 的名称、占位符类型与位置
  4. 主题配色：<a:clrScheme> 中的 12 个标准 Token + 其 sRGB 值
  5. 字体方案：Latin / East-Asian / Complex-Script 三族中英文映射
  6. 段落样式：各级文本占位符的字号 / 加粗 / 对齐 / 行距
  7. Logo 候选：按尺寸+位置启发式定位可能的 Logo Shape

全部结果封装为 TemplateProfile dataclass，可直接序列化为 JSON，
供 ppt_builder.py 消费。

依赖：python-pptx  lxml
安装：pip install python-pptx lxml
"""

from __future__ import annotations

import json
import re
import sys
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu


# ---------------------------------------------------------------------------
# 异常定义
# ---------------------------------------------------------------------------
class AnalysisError(Exception):
    """PPT解析异常基类"""
    pass


class FileFormatError(AnalysisError):
    """文件格式错误"""
    pass


class ThemeParseError(AnalysisError):
    """主题解析错误"""
    pass

# ---------------------------------------------------------------------------
# 命名空间常量
# ---------------------------------------------------------------------------
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

# 飞书/Office 常见字体别名映射（可按需扩展）
FONT_ALIAS: dict[str, str] = {
    "+mj-lt": "Calibri",
    "+mn-lt": "Calibri",
    "+mj-ea": "思源黑体",
    "+mn-ea": "微软雅黑",
    "+mj-cs": "Times New Roman",
    "+mn-cs": "Arial",
}

# 标准主题色 Token 顺序（OOXML 规范）
THEME_COLOR_TOKENS = [
    "dk1", "lt1", "dk2", "lt2",
    "accent1", "accent2", "accent3", "accent4", "accent5", "accent6",
    "hlink", "folHlink",
]

# 段落对齐枚举 → 可读字符串
ALIGN_LABEL = {
    PP_ALIGN.LEFT: "left",
    PP_ALIGN.CENTER: "center",
    PP_ALIGN.RIGHT: "right",
    PP_ALIGN.JUSTIFY: "justify",
    PP_ALIGN.DISTRIBUTE: "distribute",
    None: "inherit",
}


# ---------------------------------------------------------------------------
# 数据结构
# ---------------------------------------------------------------------------

@dataclass
class PageGeometry:
    """页面尺寸，EMU 和厘米均保存，避免调用方重复换算。"""
    width_emu: int = 0
    height_emu: int = 0
    width_cm: float = 0.0
    height_cm: float = 0.0
    orientation: str = "landscape"   # landscape | portrait | square


@dataclass
class ThemeColor:
    """单个主题色 Token。"""
    token: str = ""        # e.g. "accent1"
    hex_color: str = ""    # e.g. "FF5500"
    description: str = ""  # 友好备注


@dataclass
class FontScheme:
    """字体方案（Latin / East-Asian / Complex）。"""
    major_latin: str = ""    # 标题族 Latin
    minor_latin: str = ""    # 正文族 Latin
    major_east_asian: str = ""
    minor_east_asian: str = ""
    major_complex: str = ""
    minor_complex: str = ""


@dataclass
class ParagraphStyle:
    """文本占位符某一级别的段落样式。"""
    level: int = 0
    font_size_pt: float | None = None
    bold: bool | None = None
    italic: bool | None = None
    font_color_hex: str | None = None
    alignment: str = "inherit"
    line_spacing_pt: float | None = None   # None 表示自动


@dataclass
class PlaceholderInfo:
    """版式或母版中单个占位符的描述。"""
    idx: int = 0
    ph_type: str = ""      # e.g. "TITLE", "BODY", "PICTURE"
    left_cm: float = 0.0
    top_cm: float = 0.0
    width_cm: float = 0.0
    height_cm: float = 0.0
    paragraph_styles: list[ParagraphStyle] = field(default_factory=list)


@dataclass
class SlideLayoutInfo:
    """单个版式描述。"""
    index: int = 0
    name: str = ""
    placeholders: list[PlaceholderInfo] = field(default_factory=list)


@dataclass
class SlideMasterInfo:
    """单个母版描述。"""
    index: int = 0
    background_hex: str | None = None   # 纯色背景的 HEX，渐变/图片背景为 None
    has_gradient_bg: bool = False
    has_picture_bg: bool = False
    decoration_shape_count: int = 0     # 非占位符形状数量（装饰元素）
    logo_candidates: list[dict[str, Any]] = field(default_factory=list)
    layouts: list[SlideLayoutInfo] = field(default_factory=list)


@dataclass
class TemplateProfile:
    """完整模板结构化描述，ppt_builder 的唯一输入。"""
    source_file: str = ""
    slide_count: int = 0
    geometry: PageGeometry = field(default_factory=PageGeometry)
    theme_colors: list[ThemeColor] = field(default_factory=list)
    font_scheme: FontScheme = field(default_factory=FontScheme)
    masters: list[SlideMasterInfo] = field(default_factory=list)
    # 所有版式的平铺列表（跨母版），方便 builder 直接索引
    all_layouts: list[SlideLayoutInfo] = field(default_factory=list)
    # 出现在幻灯片中频次最高的字体（统计分析）
    dominant_fonts: list[str] = field(default_factory=list)
    # 原始幻灯片中提取的所有 RGB 色（去重后，频次降序）
    slide_colors: list[str] = field(default_factory=list)


# ---------------------------------------------------------------------------
# 工具函数
# ---------------------------------------------------------------------------

def _emu_to_cm(emu: int) -> float:
    return round(emu / 914400 * 2.54, 4)


def _resolve_font_alias(name: str) -> str:
    """将 +mj-lt 等别名解析为可读字体名。"""
    return FONT_ALIAS.get(name, name) if name else ""


def _rgb_from_element(el: etree._Element) -> str | None:
    """
    从 <a:srgbClr>、<a:sysClr>、<a:prstClr> 等节点提取 HEX 字符串。
    返回 6 位大写 HEX，失败返回 None。
    """
    srgb = el.find(qn("a:srgbClr"))
    if srgb is not None:
        val = srgb.get("val", "")
        return val.upper() if len(val) == 6 else None

    sys_clr = el.find(qn("a:sysClr"))
    if sys_clr is not None:
        # lastClr 是渲染时缓存的实际颜色
        last = sys_clr.get("lastClr", "")
        return last.upper() if len(last) == 6 else None

    return None


def _pt_from_val(val: str | None) -> float | None:
    """OOXML 字号单位 = 百分之一磅（hundredths of a point）→ pt。"""
    if val is None:
        return None
    try:
        return round(int(val) / 100, 1)
    except ValueError:
        return None


def _line_spacing_pt(sp_el: etree._Element | None) -> float | None:
    """解析 <a:lnSpc> → pt 值；百分比行距返回 None（保留默认）。"""
    if sp_el is None:
        return None
    spcPts = sp_el.find(qn("a:spcPts"))
    if spcPts is not None:
        try:
            return round(int(spcPts.get("val", "0")) / 100, 1)
        except ValueError:
            return None
    return None  # spcPct 暂不换算，保留继承


# ---------------------------------------------------------------------------
# 解析子模块
# ---------------------------------------------------------------------------

def _parse_geometry(prs: Presentation) -> PageGeometry:
    w = prs.slide_width
    h = prs.slide_height
    w_cm = _emu_to_cm(w)
    h_cm = _emu_to_cm(h)
    if w > h:
        orientation = "landscape"
    elif h > w:
        orientation = "portrait"
    else:
        orientation = "square"
    return PageGeometry(
        width_emu=int(w),
        height_emu=int(h),
        width_cm=w_cm,
        height_cm=h_cm,
        orientation=orientation,
    )


def _parse_theme_colors(prs: Presentation) -> list[ThemeColor]:
    """
    从第一个 SlideMaster 的 theme.xml 中提取 <a:clrScheme>。
    若有多个母版且主题不同，取第一个（通常是主品牌主题）。
    """
    results: list[ThemeColor] = []
    try:
        theme_part = prs.slide_masters[0].theme
    except (IndexError, AttributeError):
        return results

    root = theme_part._element
    clr_scheme = root.find(".//" + qn("a:clrScheme"))
    if clr_scheme is None:
        return results

    token_desc = {
        "dk1": "深色1（主文字色）",
        "lt1": "浅色1（背景色）",
        "dk2": "深色2",
        "lt2": "浅色2",
        "accent1": "强调色1",
        "accent2": "强调色2",
        "accent3": "强调色3",
        "accent4": "强调色4",
        "accent5": "强调色5",
        "accent6": "强调色6",
        "hlink": "超链接色",
        "folHlink": "已访问链接色",
    }

    for token in THEME_COLOR_TOKENS:
        node = clr_scheme.find(qn(f"a:{token}"))
        if node is None:
            continue
        hex_val = _rgb_from_element(node)
        results.append(ThemeColor(
            token=token,
            hex_color=hex_val or "未知",
            description=token_desc.get(token, ""),
        ))
    return results


def _parse_font_scheme(prs: Presentation) -> FontScheme:
    """从 theme.xml 的 <a:fontScheme> 提取字体方案。"""
    fs = FontScheme()
    try:
        theme_part = prs.slide_masters[0].theme
    except (IndexError, AttributeError):
        return fs

    root = theme_part._element
    font_scheme = root.find(".//" + qn("a:fontScheme"))
    if font_scheme is None:
        return fs

    def _pick(node_name: str, attr: str) -> str:
        node = font_scheme.find(f".//{qn(node_name)}")
        if node is None:
            return ""
        latin = node.find(qn("a:latin"))
        ea = node.find(qn("a:ea"))
        cs = node.find(qn("a:cs"))
        if attr == "latin":
            return _resolve_font_alias(latin.get("typeface", "") if latin is not None else "")
        if attr == "ea":
            return _resolve_font_alias(ea.get("typeface", "") if ea is not None else "")
        if attr == "cs":
            return _resolve_font_alias(cs.get("typeface", "") if cs is not None else "")
        return ""

    fs.major_latin      = _pick("a:majorFont", "latin")
    fs.major_east_asian = _pick("a:majorFont", "ea")
    fs.major_complex    = _pick("a:majorFont", "cs")
    fs.minor_latin      = _pick("a:minorFont", "latin")
    fs.minor_east_asian = _pick("a:minorFont", "ea")
    fs.minor_complex    = _pick("a:minorFont", "cs")
    return fs


def _parse_paragraph_styles(ph_el: etree._Element) -> list[ParagraphStyle]:
    """
    从占位符 XML 节点中提取各级段落样式（<a:lvl1pPr> … <a:lvl9pPr>）。
    同时兜底读取 <a:pPr> 中的单段样式。
    """
    styles: list[ParagraphStyle] = []
    # 尝试从 txBody/lstStyle 中读取各级预设
    txBody = ph_el.find(".//" + qn("p:txBody"))
    if txBody is None:
        return styles

    lstStyle = txBody.find(qn("a:lstStyle"))
    if lstStyle is not None:
        for lvl_idx in range(1, 10):
            lvl_node = lstStyle.find(qn(f"a:lvl{lvl_idx}pPr"))
            if lvl_node is None:
                continue
            ps = ParagraphStyle(level=lvl_idx - 1)
            # 对齐
            algn = lvl_node.get("algn")
            ps.alignment = algn if algn else "inherit"
            # 行距
            lnSpc = lvl_node.find(qn("a:lnSpc"))
            ps.line_spacing_pt = _line_spacing_pt(lnSpc)
            # 字体属性从第一个 <a:r>/<a:rPr> 读取
            rPr = lvl_node.find(".//" + qn("a:rPr"))
            if rPr is not None:
                ps.bold   = rPr.get("b") == "1"
                ps.italic = rPr.get("i") == "1"
                ps.font_size_pt = _pt_from_val(rPr.get("sz"))
                solid_fill = rPr.find(qn("a:solidFill"))
                if solid_fill is not None:
                    ps.font_color_hex = _rgb_from_element(solid_fill)
            styles.append(ps)

    # 若 lstStyle 没有内容，降级从实际段落提取第一段样式
    if not styles:
        for para in txBody.findall(qn("a:p")):
            pPr = para.find(qn("a:pPr"))
            rPr = para.find(".//" + qn("a:rPr"))
            ps = ParagraphStyle(level=0)
            if pPr is not None:
                algn = pPr.get("algn")
                ps.alignment = algn or "inherit"
                lnSpc = pPr.find(qn("a:lnSpc"))
                ps.line_spacing_pt = _line_spacing_pt(lnSpc)
            if rPr is not None:
                ps.bold = rPr.get("b") == "1"
                ps.italic = rPr.get("i") == "1"
                ps.font_size_pt = _pt_from_val(rPr.get("sz"))
            styles.append(ps)
            break  # 只取第一段

    return styles


def _parse_placeholder(sp: Any) -> PlaceholderInfo:
    """将 pptx Placeholder Shape 转为 PlaceholderInfo。"""
    ph_type = str(sp.placeholder_format.type).replace("PP_PLACEHOLDER.", "")
    info = PlaceholderInfo(
        idx=sp.placeholder_format.idx,
        ph_type=ph_type,
        left_cm=_emu_to_cm(sp.left or 0),
        top_cm=_emu_to_cm(sp.top or 0),
        width_cm=_emu_to_cm(sp.width or 0),
        height_cm=_emu_to_cm(sp.height or 0),
        paragraph_styles=_parse_paragraph_styles(sp._element),
    )
    return info


def _is_logo_candidate(shape: Any, page_w_emu: int, page_h_emu: int) -> bool:
    """
    启发式判断 Shape 是否为 Logo：
      - 是图片类型（Picture / OLE）
      - 尺寸较小（宽/高均 < 页面的 30%）
      - 位于四个角落（距任一角 < 页面尺寸的 20%）
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    _LOGO_TYPES = {
        MSO_SHAPE_TYPE.PICTURE,
        MSO_SHAPE_TYPE.LINKED_PICTURE,
        MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,   # python-pptx 1.0+
        MSO_SHAPE_TYPE.LINKED_OLE_OBJECT,     # python-pptx 1.0+
    }
    if shape.shape_type not in _LOGO_TYPES:
        return False
    if shape.width is None or shape.height is None:
        return False
    w_ratio = shape.width / page_w_emu
    h_ratio = shape.height / page_h_emu
    if w_ratio > 0.3 or h_ratio > 0.3:
        return False
    # 检查是否在角落
    corners = [
        (0, 0),
        (page_w_emu - shape.width, 0),
        (0, page_h_emu - shape.height),
        (page_w_emu - shape.width, page_h_emu - shape.height),
    ]
    for cx, cy in corners:
        if (
            abs((shape.left or 0) - cx) < page_w_emu * 0.20
            and abs((shape.top or 0) - cy) < page_h_emu * 0.20
        ):
            return True
    return False


def _get_bg_hex(master_or_layout: Any) -> tuple[str | None, bool, bool]:
    """
    提取背景色。返回 (hex | None, has_gradient, has_picture)。
    """
    bg = master_or_layout.background
    fill = bg.fill
    try:
        if fill.type is None:
            return None, False, False
        from pptx.enum.dml import MSO_THEME_COLOR
        from pptx.oxml.ns import qn as _qn
        fill_xml = fill._xPr
        # 纯色
        solidFill = fill_xml.find(_qn("a:solidFill"))
        if solidFill is not None:
            hex_val = _rgb_from_element(solidFill)
            return hex_val, False, False
        # 渐变
        gradFill = fill_xml.find(_qn("a:gradFill"))
        if gradFill is not None:
            return None, True, False
        # 图片
        blipFill = fill_xml.find(_qn("a:blipFill"))
        if blipFill is not None:
            return None, False, True
    except Exception:
        pass
    return None, False, False


def _parse_master(master: Any, idx: int, page_w: int, page_h: int) -> SlideMasterInfo:
    """解析单个 SlideMaster。"""
    bg_hex, has_grad, has_pic = _get_bg_hex(master)
    deco_count = sum(
        1 for sp in master.shapes
        if not sp.is_placeholder
    )
    logo_candidates = []
    for sp in master.shapes:
        if _is_logo_candidate(sp, page_w, page_h):
            logo_candidates.append({
                "name": sp.name,
                "left_cm": _emu_to_cm(sp.left or 0),
                "top_cm": _emu_to_cm(sp.top or 0),
                "width_cm": _emu_to_cm(sp.width or 0),
                "height_cm": _emu_to_cm(sp.height or 0),
            })

    layouts: list[SlideLayoutInfo] = []
    for l_idx, layout in enumerate(master.slide_layouts):
        l_info = SlideLayoutInfo(index=l_idx, name=layout.name)
        for sp in layout.placeholders:
            try:
                l_info.placeholders.append(_parse_placeholder(sp))
            except Exception:
                pass
        layouts.append(l_info)

    return SlideMasterInfo(
        index=idx,
        background_hex=bg_hex,
        has_gradient_bg=has_grad,
        has_picture_bg=has_pic,
        decoration_shape_count=deco_count,
        logo_candidates=logo_candidates,
        layouts=layouts,
    )


def _collect_slide_colors(prs: Presentation) -> list[str]:
    """
    遍历所有幻灯片的 Shape，收集显式 RGB 色，按频次降序去重返回。
    采样最多 10 张幻灯片以控制耗时。
    注意：python-pptx 1.0+ 的 SlidePart 不支持切片，需先转为 list。
    """
    from collections import Counter
    counter: Counter = Counter()
    sample = list(prs.slides)[:10]   # 修复：先转 list 再切片
    for slide in sample:
        for shape in slide.shapes:
            el = shape._element
            for srgb in el.iter(qn("a:srgbClr")):
                val = srgb.get("val", "").upper()
                if re.fullmatch(r"[0-9A-F]{6}", val):
                    counter[val] += 1
    return [color for color, _ in counter.most_common(20)]


def _collect_dominant_fonts(prs: Presentation) -> list[str]:
    """
    遍历所有幻灯片文本，统计出现最多的字体名（去别名后）。
    注意：python-pptx 1.0+ 的 SlidePart 不支持切片，需先转为 list。
    """
    from collections import Counter
    counter: Counter = Counter()
    sample = list(prs.slides)[:10]   # 修复：先转 list 再切片
    for slide in sample:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    fn = run.font.name
                    if fn:
                        counter[_resolve_font_alias(fn)] += 1
    return [font for font, _ in counter.most_common(10)]


# ---------------------------------------------------------------------------
# 主入口
# ---------------------------------------------------------------------------

def analyze(pptx_path: str | Path) -> TemplateProfile:
    """
    解析 .pptx 文件，返回 TemplateProfile。

    Parameters
    ----------
    pptx_path : str | Path
        本地 .pptx 文件路径。

    Returns
    -------
    TemplateProfile
        结构化模板描述对象。

    Raises
    ------
    FileNotFoundError
        文件不存在时抛出。
    ValueError
        文件扩展名不是 .pptx 时抛出（拒绝旧 .ppt 格式）。
    RuntimeError
        文件损坏或无法解析时抛出（包装底层异常）。
    """
    path = Path(pptx_path)

    # --- 前置校验 ---
    if not path.exists():
        raise FileNotFoundError(f"文件不存在：{path}")
    if path.suffix.lower() != ".pptx":
        raise ValueError(
            f"仅支持 .pptx 格式，当前文件为 {path.suffix}。"
            "请用 PowerPoint / WPS 另存为 .pptx 后重试。"
        )

    # --- 加载 ---
    try:
        prs = Presentation(str(path))
    except Exception as exc:
        raise RuntimeError(f"PPT 文件加载失败，可能已损坏：{exc}") from exc

    profile = TemplateProfile(source_file=str(path))

    # 1. 页面几何
    profile.geometry = _parse_geometry(prs)
    profile.slide_count = len(prs.slides)

    # 2. 主题配色
    profile.theme_colors = _parse_theme_colors(prs)

    # 3. 字体方案
    profile.font_scheme = _parse_font_scheme(prs)

    # 4. 母版 + 版式
    page_w = int(prs.slide_width)
    page_h = int(prs.slide_height)
    for m_idx, master in enumerate(prs.slide_masters):
        m_info = _parse_master(master, m_idx, page_w, page_h)
        profile.masters.append(m_info)
        profile.all_layouts.extend(m_info.layouts)

    # 5. 幻灯片级颜色统计（补充主题色未覆盖的实际用色）
    profile.slide_colors = _collect_slide_colors(prs)

    # 6. 主导字体统计
    profile.dominant_fonts = _collect_dominant_fonts(prs)

    return profile


def analyze_to_json(pptx_path: str | Path, indent: int = 2) -> str:
    """解析并返回 JSON 字符串，方便调试和跨模块传递。"""
    profile = analyze(pptx_path)
    return json.dumps(asdict(profile), ensure_ascii=False, indent=indent)


# ---------------------------------------------------------------------------
# CLI 自验证入口（python ppt_analyzer.py your_file.pptx）
# ---------------------------------------------------------------------------

def _print_summary(profile: TemplateProfile) -> None:
    """控制台打印人类可读摘要。"""
    geo = profile.geometry
    print("=" * 60)
    print(f"  PPT 解析报告：{Path(profile.source_file).name}")
    print("=" * 60)
    print(f"  页面尺寸  : {geo.width_cm} cm × {geo.height_cm} cm ({geo.orientation})")
    print(f"  幻灯片数  : {profile.slide_count} 页")
    print(f"  母版数    : {len(profile.masters)} 个")
    print(f"  版式数    : {len(profile.all_layouts)} 种")

    print("\n  【主题配色】")
    for tc in profile.theme_colors:
        hex_color = tc.hex_color if tc.hex_color and tc.hex_color != "未知" else "000000"
        if len(hex_color) == 6:
            try:
                bar = f"\033[48;2;{int(hex_color[0:2],16)};{int(hex_color[2:4],16)};{int(hex_color[4:6],16)}m  \033[0m"
            except ValueError:
                bar = "  "
        else:
            bar = "  "
        print(f"    {bar} #{hex_color}  {tc.token:10s}  {tc.description}")

    print("\n  【字体方案】")
    fs = profile.font_scheme
    print(f"    标题-Latin   : {fs.major_latin}")
    print(f"    正文-Latin   : {fs.minor_latin}")
    print(f"    标题-东亚    : {fs.major_east_asian}")
    print(f"    正文-东亚    : {fs.minor_east_asian}")

    print("\n  【版式清单】")
    for layout in profile.all_layouts:
        ph_names = ", ".join(p.ph_type for p in layout.placeholders) or "无占位符"
        print(f"    [{layout.index:02d}] {layout.name:<30s}  占位符: {ph_names}")

    if profile.dominant_fonts:
        print(f"\n  【幻灯片主导字体（Top5）】")
        for f in profile.dominant_fonts[:5]:
            print(f"    · {f}")

    if profile.slide_colors:
        print(f"\n  【幻灯片实际用色（Top10）】")
        row = "    " + "  ".join(f"#{c}" for c in profile.slide_colors[:10])
        print(row)

    logo_count = sum(len(m.logo_candidates) for m in profile.masters)
    if logo_count:
        print(f"\n  【Logo 候选位置】 共检测到 {logo_count} 个")
        for m in profile.masters:
            for lc in m.logo_candidates:
                print(f"    母版[{m.index}] {lc['name']} "
                      f"@ ({lc['left_cm']:.1f}, {lc['top_cm']:.1f}) cm  "
                      f"{lc['width_cm']:.1f}×{lc['height_cm']:.1f} cm")
    print("=" * 60)


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法：python ppt_analyzer.py <your_file.pptx> [--json]")
        print("      --json  输出完整 JSON，否则输出人类可读摘要")
        sys.exit(1)

    target = sys.argv[1]
    output_json = "--json" in sys.argv

    try:
        if output_json:
            print(analyze_to_json(target))
        else:
            profile = analyze(target)
            _print_summary(profile)
    except (FileNotFoundError, ValueError, RuntimeError) as e:
        print(f"[错误] {e}", file=sys.stderr)
        sys.exit(2)
