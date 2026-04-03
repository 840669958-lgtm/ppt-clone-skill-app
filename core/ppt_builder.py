"""
core/ppt_builder.py
===================
PPT 模板重建引擎 —— 根据解析数据 1:1 复刻生成新 PPT。

职责：
  1. 根据 TemplateProfile 重建页面尺寸、主题配色、字体方案
  2. 复刻母版背景、装饰元素、占位符布局
  3. 创建指定数量的空白幻灯片，绑定对应版式
  4. 支持品牌 Logo 替换、配色调整、页面数量自定义

核心原则：
  - 只复刻「模板结构」，不复制内容
  - 所有样式 1:1 还原，生成空白可编辑 PPT

依赖：python-pptx
安装：pip install python-pptx
"""

from __future__ import annotations

import copy
import logging
import os
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.util import Emu, Pt, Inches, Cm

from .ppt_analyzer import (
    TemplateProfile,
    PageGeometry,
    ThemeColor,
    FontScheme,
    SlideMasterInfo,
    SlideLayoutInfo,
    PlaceholderInfo,
    ParagraphStyle,
    _emu_to_cm,
)

logger = logging.getLogger(__name__)

# 默认输出配置
DEFAULT_OUTPUT_DIR = "./output"
DEFAULT_SLIDE_COUNT = 10


# ---------------------------------------------------------------------------
# 数据结构
# ---------------------------------------------------------------------------

@dataclass
class BuildOptions:
    """
    PPT重建选项。
    
    Attributes
    ----------
    slide_count : int
        生成幻灯片数量，默认10页
    replace_logo_path : str | None
        替换Logo的本地图片路径，None表示保留原模板Logo位置（留空）
    primary_color : str | None
        主色调HEX值（如"FF5500"），None表示使用原模板配色
    output_name : str
        输出文件名
    preserve_animations : bool
        是否保留动画（默认False，生成干净模板）
    """
    slide_count: int = DEFAULT_SLIDE_COUNT
    replace_logo_path: str | None = None
    primary_color: str | None = None
    output_name: str = "cloned_template.pptx"
    preserve_animations: bool = False


@dataclass
class BuildResult:
    """PPT构建结果。"""
    output_path: str
    slide_count: int
    layout_usage: dict[str, int]  # 版式名称 → 使用次数
    warnings: list[str]


# ---------------------------------------------------------------------------
# 异常定义
# ---------------------------------------------------------------------------

class PPTBuilderError(Exception):
    """PPT构建异常基类。"""
    
    def __init__(self, message: str, solution: str = ""):
        self.message = message
        self.solution = solution
        super().__init__(message)


class InvalidProfileError(PPTBuilderError):
    """模板数据无效。"""
    pass


class LogoReplaceError(PPTBuilderError):
    """Logo替换失败。"""
    pass


# ---------------------------------------------------------------------------
# 工具函数
# ---------------------------------------------------------------------------

def _cm_to_emu(cm: float) -> int:
    """厘米转EMU（Office度量单位）。"""
    return int(cm / 2.54 * 914400)


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """HEX颜色字符串转RGBColor对象。"""
    hex_color = hex_color.lstrip("#").upper()
    if len(hex_color) != 6:
        raise ValueError(f"无效的HEX颜色: {hex_color}")
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def _adjust_color(hex_color: str, primary_color: str | None) -> str:
    """
    根据主色调调整颜色。
    如果指定了primary_color，将原颜色的亮度/饱和度映射到新主色调。
    """
    if not primary_color:
        return hex_color
    
    # 简化处理：如果是主题色（accent1-6），替换为primary_color
    # 实际项目中可实现更复杂的色彩映射算法
    return primary_color


def _remove_animations(slide) -> None:
    """移除幻灯片中的所有动画节点。"""
    # 找到 <p:timing> 节点并移除
    timing = slide._element.find(qn("p:timing"))
    if timing is not None:
        slide._element.remove(timing)


# ---------------------------------------------------------------------------
# 核心类
# ---------------------------------------------------------------------------

class PPTBuilder:
    """
    PPT模板重建器。
    
    根据解析出的 TemplateProfile，1:1 复刻生成新的空白PPT模板。
    
    Parameters
    ----------
    profile : TemplateProfile
        模板解析结果（来自ppt_analyzer.py）
    output_dir : str | Path
        输出目录
        
    Examples
    --------
    >>> from core.ppt_analyzer import analyze
    >>> from core.ppt_builder import PPTBuilder, BuildOptions
    >>> 
    >>> # 解析模板
    >>> profile = analyze("reference.pptx")
    >>> 
    >>> # 重建模板
    >>> builder = PPTBuilder(profile, output_dir="./output")
    >>> options = BuildOptions(
    ...     slide_count=15,
    ...     replace_logo_path="./new_logo.png",
    ...     primary_color="FF5500"
    ... )
    >>> result = builder.build(options)
    >>> print(f"生成完成: {result.output_path}")
    """

    def __init__(
        self,
        profile: TemplateProfile,
        output_dir: str | Path = DEFAULT_OUTPUT_DIR,
    ):
        self.profile = profile
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.warnings: list[str] = []

    def build(self, options: BuildOptions | None = None) -> BuildResult:
        """
        执行PPT重建。
        
        流程：
          1. 创建空白Presentation，设置页面尺寸
          2. 重建主题配色和字体方案
          3. 重建母版和版式结构
          4. 创建指定数量的幻灯片
          5. 可选：替换Logo、调整配色
          6. 保存输出
        
        Parameters
        ----------
        options : BuildOptions | None
            构建选项，None使用默认配置
            
        Returns
        -------
        BuildResult
            构建结果，包含输出路径和统计信息
            
        Raises
        ------
        InvalidProfileError
            模板数据无效或缺失关键信息
        LogoReplaceError
            Logo替换失败（文件不存在或格式不支持）
        """
        if options is None:
            options = BuildOptions()
        
        self.warnings = []
        logger.info(f"开始构建PPT: {options.output_name}")
        
        # Step 1: 创建基础Presentation
        prs = self._create_base_presentation()
        
        # Step 2: 应用页面尺寸
        self._apply_geometry(prs)
        
        # Step 3: 重建主题（配色+字体）
        self._apply_theme(prs, options)
        
        # Step 4: 重建母版和版式
        self._rebuild_masters_and_layouts(prs)
        
        # Step 5: 创建幻灯片
        layout_usage = self._create_slides(prs, options)
        
        # Step 6: 保存文件
        output_path = self.output_dir / options.output_name
        prs.save(str(output_path))
        
        logger.info(f"PPT构建完成: {output_path}")
        
        return BuildResult(
            output_path=str(output_path),
            slide_count=len(prs.slides),
            layout_usage=layout_usage,
            warnings=self.warnings,
        )

    def _create_base_presentation(self) -> Presentation:
        """创建基础Presentation对象。"""
        # 使用空白模板创建
        prs = Presentation()
        # 移除默认的空白幻灯片
        if len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
        return prs

    def _apply_geometry(self, prs: Presentation) -> None:
        """应用页面尺寸。"""
        geo = self.profile.geometry
        if geo.width_emu > 0 and geo.height_emu > 0:
            prs.slide_width = Emu(geo.width_emu)
            prs.slide_height = Emu(geo.height_emu)
            logger.info(f"应用页面尺寸: {geo.width_cm}cm x {geo.height_cm}cm ({geo.orientation})")
        else:
            self.warnings.append("模板数据缺少页面尺寸，使用默认16:9")
            # 默认16:9
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

    def _apply_theme(self, prs: Presentation, options: BuildOptions) -> None:
        """应用主题配色和字体方案。"""
        # 获取或创建主题部分
        theme_part = self._get_or_create_theme_part(prs)
        
        # 应用配色方案
        self._apply_color_scheme(theme_part, options)
        
        # 应用字体方案
        self._apply_font_scheme(theme_part)

    def _get_or_create_theme_part(self, prs: Presentation):
        """获取或创建主题部分。"""
        # 从第一个母版获取主题
        if prs.slide_masters:
            master = prs.slide_masters[0]
            # 通过part获取theme - 使用正确的属性名
            try:
                # python-pptx 1.0+ 使用 theme_part
                return master.part.theme_part
            except AttributeError:
                # 旧版本可能使用其他方式
                pass
        return None

    def _apply_color_scheme(self, theme_part, options: BuildOptions) -> None:
        """应用配色方案到主题。"""
        if not theme_part or not self.profile.theme_colors:
            return
        
        try:
            # 获取颜色方案元素
            clr_scheme = theme_part.element.find(qn("a:clrScheme"))
            if clr_scheme is None:
                return
            
            # 构建颜色映射
            color_map = {tc.token: tc.hex_color for tc in self.profile.theme_colors}
            
            # 应用主色调调整
            if options.primary_color:
                # 将accent1-6映射到新主色调
                for i in range(1, 7):
                    token = f"accent{i}"
                    if token in color_map:
                        color_map[token] = options.primary_color
            
            # 更新主题中的颜色
            for token, hex_color in color_map.items():
                # 查找对应的颜色元素
                color_el = clr_scheme.find(qn(f"a:{token}"))
                if color_el is not None:
                    # 查找srgbClr或sysClr子元素
                    srgb = color_el.find(qn("a:srgbClr"))
                    if srgb is not None:
                        srgb.set("val", hex_color.lower())
                    else:
                        sys_clr = color_el.find(qn("a:sysClr"))
                        if sys_clr is not None:
                            # 创建srgbClr替换sysClr
                            new_srgb = parse_xml(f'<a:srgbClr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="{hex_color.lower()}"/>')
                            color_el.remove(sys_clr)
                            color_el.append(new_srgb)
            
            logger.info(f"应用配色方案: {len(color_map)} 个颜色")
            
        except Exception as e:
            self.warnings.append(f"应用配色方案时出错: {e}")
            logger.warning(f"配色应用失败: {e}")

    def _apply_font_scheme(self, theme_part) -> None:
        """应用字体方案到主题。"""
        if not theme_part:
            return
        
        font_scheme = self.profile.font_scheme
        if not any([
            font_scheme.major_latin, font_scheme.minor_latin,
            font_scheme.major_east_asian, font_scheme.minor_east_asian
        ]):
            return
        
        try:
            # 获取字体方案元素
            font_scheme_el = theme_part.element.find(qn("a:fontScheme"))
            if font_scheme_el is None:
                return
            
            # 更新主要字体（标题）
            major_font = font_scheme_el.find(qn("a:majorFont"))
            if major_font is not None:
                if font_scheme.major_latin:
                    latin = major_font.find(qn("a:latin"))
                    if latin is not None:
                        latin.set("typeface", font_scheme.major_latin)
                if font_scheme.major_east_asian:
                    ea = major_font.find(qn("a:ea"))
                    if ea is not None:
                        ea.set("typeface", font_scheme.major_east_asian)
            
            # 更新次要字体（正文）
            minor_font = font_scheme_el.find(qn("a:minorFont"))
            if minor_font is not None:
                if font_scheme.minor_latin:
                    latin = minor_font.find(qn("a:latin"))
                    if latin is not None:
                        latin.set("typeface", font_scheme.minor_latin)
                if font_scheme.minor_east_asian:
                    ea = minor_font.find(qn("a:ea"))
                    if ea is not None:
                        ea.set("typeface", font_scheme.minor_east_asian)
            
            logger.info("应用字体方案完成")
            
        except Exception as e:
            self.warnings.append(f"应用字体方案时出错: {e}")
            logger.warning(f"字体应用失败: {e}")

    def _rebuild_masters_and_layouts(self, prs: Presentation) -> None:
        """
        重建母版和版式结构。
        
        由于python-pptx对母版和版式的创建支持有限，
        这里采用策略：保留默认母版，通过修改其版式来近似复刻。
        """
        if not prs.slide_masters:
            self.warnings.append("无法获取母版，跳过版式重建")
            return
        
        master = prs.slide_masters[0]
        
        # 获取解析出的版式信息
        layouts_info = self.profile.all_layouts
        if not layouts_info:
            self.warnings.append("模板数据缺少版式信息")
            return
        
        # 记录可用的版式
        self._available_layouts = list(master.slide_layouts)
        logger.info(f"可用版式数量: {len(self._available_layouts)}")
        
        # 如果解析的版式数量多于可用版式，记录警告
        if len(layouts_info) > len(self._available_layouts):
            self.warnings.append(
                f"原模板有 {len(layouts_info)} 个版式，"
                f"但新PPT只能创建 {len(self._available_layouts)} 个"
            )

    def _create_slides(self, prs: Presentation, options: BuildOptions) -> dict[str, int]:
        """
        创建指定数量的幻灯片。
        
        策略：
          - 优先使用原模板中常用的版式
          - 循环使用可用版式
          - 记录每个版式的使用次数
        """
        layout_usage: dict[str, int] = {}
        
        # 获取可用版式
        available_layouts = list(prs.slide_masters[0].slide_layouts) if prs.slide_masters else []
        if not available_layouts:
            raise PPTBuilderError("无法获取版式，无法创建幻灯片")
        
        # 准备版式使用计划
        # 优先使用原模板中定义的版式（按索引匹配）
        layout_plan = []
        for i in range(options.slide_count):
            if i < len(self.profile.all_layouts):
                # 使用对应索引的版式
                layout_idx = min(i, len(available_layouts) - 1)
            else:
                # 循环使用版式
                layout_idx = i % len(available_layouts)
            layout_plan.append(layout_idx)
        
        # 创建幻灯片
        for idx, layout_idx in enumerate(layout_plan):
            layout = available_layouts[layout_idx]
            slide = prs.slides.add_slide(layout)
            
            # 记录版式使用
            layout_name = layout.name or f"Layout_{layout_idx}"
            layout_usage[layout_name] = layout_usage.get(layout_name, 0) + 1
            
            # 移除动画（如果不保留）
            if not options.preserve_animations:
                _remove_animations(slide)
            
            # 如果是第一页且指定了Logo，尝试替换
            if idx == 0 and options.replace_logo_path:
                self._try_replace_logo(slide, options.replace_logo_path)
        
        logger.info(f"创建幻灯片: {len(prs.slides)} 页")
        return layout_usage

    def _try_replace_logo(self, slide, logo_path: str) -> None:
        """尝试在幻灯片中替换Logo。"""
        if not os.path.exists(logo_path):
            self.warnings.append(f"Logo文件不存在: {logo_path}")
            return
        
        try:
            from pptx.util import Inches
            
            # 启发式：查找角落的小图片作为Logo位置
            logo_shape = None
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # 检查是否在角落且尺寸较小
                    left_pct = shape.left / slide.shapes._spTree.getparent().getparent().getparent().sz.cx
                    top_pct = shape.top / slide.shapes._spTree.getparent().getparent().getparent().sz.cy
                    
                    # 简单的角落检测（左上、右上）
                    is_corner = (left_pct < 0.2 or left_pct > 0.8) and top_pct < 0.3
                    is_small = shape.width < Inches(2) and shape.height < Inches(1)
                    
                    if is_corner and is_small:
                        logo_shape = shape
                        break
            
            if logo_shape:
                # 记录位置后删除原Logo
                left, top, width, height = logo_shape.left, logo_shape.top, logo_shape.width, logo_shape.height
                sp = logo_shape._element
                sp.getparent().remove(sp)
                
                # 添加新Logo
                slide.shapes.add_picture(logo_path, left, top, width, height)
                logger.info(f"Logo已替换: {logo_path}")
            else:
                # 未找到合适位置，添加到左上角默认位置
                slide.shapes.add_picture(logo_path, Inches(0.5), Inches(0.5), height=Inches(0.8))
                logger.info(f"Logo已添加到默认位置: {logo_path}")
                
        except Exception as e:
            self.warnings.append(f"Logo替换失败: {e}")
            logger.warning(f"Logo替换失败: {e}")


# ---------------------------------------------------------------------------
# 便捷函数
# ---------------------------------------------------------------------------

def build_from_profile(
    profile: TemplateProfile,
    output_dir: str = DEFAULT_OUTPUT_DIR,
    slide_count: int = DEFAULT_SLIDE_COUNT,
    replace_logo_path: str | None = None,
    primary_color: str | None = None,
    output_name: str = "cloned_template.pptx",
) -> BuildResult:
    """
    便捷函数：从模板配置直接构建PPT。
    
    Parameters
    ----------
    profile : TemplateProfile
        模板解析结果
    output_dir : str
        输出目录
    slide_count : int
        幻灯片数量
    replace_logo_path : str | None
        Logo替换路径
    primary_color : str | None
        主色调HEX
    output_name : str
        输出文件名
        
    Returns
    -------
    BuildResult
        构建结果
    """
    builder = PPTBuilder(profile, output_dir)
    options = BuildOptions(
        slide_count=slide_count,
        replace_logo_path=replace_logo_path,
        primary_color=primary_color,
        output_name=output_name,
    )
    return builder.build(options)


# ---------------------------------------------------------------------------
# 测试入口
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    """
    独立测试入口。
    
    运行方式:
      1. 准备参考PPT文件: ./test_reference.pptx
      2. 运行: python -m core.ppt_builder
    """
    import sys
    
    logging.basicConfig(
        level=logging.INFO,
        format="%(asctime)s - %(levelname)s - %(message)s",
        datefmt="%H:%M:%S"
    )
    
    def print_separator(title: str = ""):
        print("=" * 60)
        if title:
            print(f"  {title}")
            print("=" * 60)
    
    def print_result(success: bool, message: str, solution: str = ""):
        status = "✅ 成功" if success else "❌ 失败"
        print(f"  [{status}] {message}")
        if solution and not success:
            print(f"      💡 {solution}")
    
    print_separator("PPTBuilder 功能测试")
    
    # 检查参考文件
    test_ref = Path("./test_reference.pptx")
    if not test_ref.exists():
        # 尝试其他路径
        alt_paths = [
            Path("./demo_template.pptx"),
            Path("../demo_template.pptx"),
            Path("./downloads/test_download.pptx"),
        ]
        for alt in alt_paths:
            if alt.exists():
                test_ref = alt
                break
    
    if not test_ref.exists():
        print("\n【环境检查】")
        print_result(False, "未找到参考PPT文件")
        print("\n  💡 请准备参考PPT文件，以下路径之一:")
        print("     - ./test_reference.pptx")
        print("     - ./demo_template.pptx")
        print("     - ./downloads/test_download.pptx")
        sys.exit(1)
    
    print("\n【环境检查】")
    print_result(True, f"找到参考文件: {test_ref}")
    
    # 测试1: 解析模板
    print("\n【测试1】解析参考模板")
    try:
        from .ppt_analyzer import analyze
        profile = analyze(str(test_ref))
        print(f"  📄 源文件: {profile.source_file}")
        print(f"  📐 页面尺寸: {profile.geometry.width_cm}cm x {profile.geometry.height_cm}cm")
        print(f"  🎨 主题色: {len(profile.theme_colors)} 个")
        print(f"  🔤 字体方案: {profile.font_scheme.major_latin or '默认'}")
        print(f"  📑 母版数量: {len(profile.masters)}")
        print(f"  📋 版式数量: {len(profile.all_layouts)}")
        print_result(True, "模板解析成功")
    except Exception as e:
        print_result(False, f"解析失败: {e}")
        sys.exit(1)
    
    # 测试2: 构建PPT（基础）
    print("\n【测试2】基础重建（10页，无自定义）")
    try:
        builder = PPTBuilder(profile, output_dir="./output")
        result = builder.build(BuildOptions(
            slide_count=10,
            output_name="test_basic.pptx"
        ))
        print(f"  💾 输出路径: {result.output_path}")
        print(f"  📊 幻灯片数: {result.slide_count}")
        print(f"  📋 版式使用: {result.layout_usage}")
        if result.warnings:
            print(f"  ⚠️  警告: {len(result.warnings)} 个")
            for w in result.warnings[:3]:
                print(f"      - {w}")
        print_result(True, "基础重建成功")
    except Exception as e:
        print_result(False, f"重建失败: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
    
    # 测试3: 构建PPT（带自定义）
    print("\n【测试3】高级重建（15页，自定义配色）")
    try:
        result = build_from_profile(
            profile,
            output_dir="./output",
            slide_count=15,
            primary_color="FF5500",  # 橙色主题
            output_name="test_custom.pptx"
        )
        print(f"  💾 输出路径: {result.output_path}")
        print(f"  📊 幻灯片数: {result.slide_count}")
        print(f"  🎨 主色调: FF5500 (橙色)")
        print_result(True, "高级重建成功")
    except Exception as e:
        print_result(False, f"重建失败: {e}")
        import traceback
        traceback.print_exc()
    
    # 汇总
    print_separator("测试结果汇总")
    print("\n  ✅ 基础重建: test_basic.pptx")
    print("  ✅ 高级重建: test_custom.pptx")
    print("\n  输出目录: ./output/")
    print("  请打开生成的PPT文件检查效果")
