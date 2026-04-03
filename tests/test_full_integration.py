# -*- coding: utf-8 -*-
"""
main.py 本地全链路集成测试
==========================

完整模拟用户真实使用流程，覆盖全链路：
【用户输入】→ 鉴权 → 解析链接下载PPT → 模板解析 → 复刻生成新PPT → 上传飞书 → 输出新PPT飞书链接

运行方式:
    # 1. 设置环境变量（真实飞书环境测试）
    set FEISHU_APP_ID=cli_xxx
    set FEISHU_APP_SECRET=xxx
    python tests/test_full_integration.py --mode real
    
    # 2. 本地模拟测试（无需飞书环境）
    python tests/test_full_integration.py --mode mock
    
    # 3. 混合模式（使用本地文件，跳过下载）
    python tests/test_full_integration.py --mode hybrid --input ./test_files/template.pptx

测试报告:
    测试完成后自动生成报告，包含：
    - 每一步执行状态
    - 总耗时统计
    - 成功率分析
    - 优化建议
"""

import os
import sys
import time
import json
import argparse
import traceback
import tempfile
from pathlib import Path
from datetime import datetime
from dataclasses import dataclass, field, asdict
from typing import Optional, Dict, Any, List, Callable
from unittest.mock import Mock, patch, MagicMock

# 添加项目根目录到路径
project_root = Path(__file__).parent.parent
sys.path.insert(0, str(project_root))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# 导入被测模块
from main import AppConfig, PPTCloneWorkflow, main as main_entry
from feishu.message import PPTShareMessage, CloneParameters, extract_ppt_urls
from feishu.auth import FeishuAuthError
from core.ppt_analyzer import analyze
from core.ppt_builder import PPTBuilder, BuildOptions


# ==================== 测试配置 ====================

@dataclass
class TestConfig:
    """测试配置"""
    mode: str = "mock"  # mock, real, hybrid
    test_ppt_url: str = "https://test.feishu.cn/slides/sldcnTest123"
    local_input_path: Optional[str] = None
    output_dir: str = "./output/test_integration"
    slide_count: int = 5
    primary_color: Optional[str] = None
    verbose: bool = True


@dataclass
class TestStep:
    """测试步骤结果"""
    name: str
    status: str  # success, failed, skipped
    duration: float
    message: str = ""
    details: Dict[str, Any] = field(default_factory=dict)
    error: Optional[str] = None
    suggestion: Optional[str] = None


@dataclass
class TestReport:
    """完整测试报告"""
    start_time: datetime
    end_time: Optional[datetime] = None
    total_duration: float = 0.0
    steps: List[TestStep] = field(default_factory=list)
    summary: Dict[str, Any] = field(default_factory=dict)
    
    def add_step(self, step: TestStep):
        self.steps.append(step)
    
    def finalize(self):
        self.end_time = datetime.now()
        self.total_duration = sum(s.duration for s in self.steps)
        
        success_count = sum(1 for s in self.steps if s.status == "success")
        failed_count = sum(1 for s in self.steps if s.status == "failed")
        skipped_count = sum(1 for s in self.steps if s.status == "skipped")
        
        self.summary = {
            "total_steps": len(self.steps),
            "success": success_count,
            "failed": failed_count,
            "skipped": skipped_count,
            "success_rate": f"{success_count / len(self.steps) * 100:.1f}%" if self.steps else "0%",
            "total_duration": f"{self.total_duration:.2f}s"
        }
    
    def to_dict(self) -> Dict[str, Any]:
        return {
            "start_time": self.start_time.isoformat(),
            "end_time": self.end_time.isoformat() if self.end_time else None,
            "total_duration": self.total_duration,
            "summary": self.summary,
            "steps": [
                {
                    "name": s.name,
                    "status": s.status,
                    "duration": f"{s.duration:.2f}s",
                    "message": s.message,
                    "details": s.details,
                    "error": s.error,
                    "suggestion": s.suggestion
                }
                for s in self.steps
            ]
        }


# ==================== 测试工具类 ====================

class TestLogger:
    """测试日志记录器"""
    
    def __init__(self, verbose: bool = True):
        self.verbose = verbose
        self.logs: List[str] = []
    
    def info(self, message: str):
        log = f"[INFO] {message}"
        self.logs.append(log)
        if self.verbose:
            print(f"  ℹ️  {message}")
    
    def success(self, message: str):
        log = f"[SUCCESS] {message}"
        self.logs.append(log)
        if self.verbose:
            print(f"  ✅ {message}")
    
    def warning(self, message: str):
        log = f"[WARNING] {message}"
        self.logs.append(log)
        if self.verbose:
            print(f"  ⚠️  {message}")
    
    def error(self, message: str):
        log = f"[ERROR] {message}"
        self.logs.append(log)
        if self.verbose:
            print(f"  ❌ {message}")
    
    def section(self, title: str):
        separator = "=" * 60
        self.logs.append(separator)
        self.logs.append(f"[SECTION] {title}")
        self.logs.append(separator)
        if self.verbose:
            print(f"\n{separator}")
            print(f"  {title}")
            print(f"{separator}")
    
    def step(self, step_num: int, total: int, title: str):
        message = f"[{step_num}/{total}] {title}"
        self.logs.append(message)
        if self.verbose:
            print(f"\n🔷 {message}")
    
    def get_logs(self) -> str:
        return "\n".join(self.logs)


class MockFeishuEnvironment:
    """模拟飞书环境"""
    
    def __init__(self, test_files_dir: Path):
        self.test_files_dir = test_files_dir
        self.mock_files: Dict[str, Path] = {}
        self._create_mock_files()
    
    def _create_mock_files(self):
        """创建模拟测试文件"""
        # 创建常规商务模板
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # 标题页
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "商务报告模板"
        slide.placeholders[1].text = "2026年度总结"
        
        # 内容页
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "核心指标"
        slide.placeholders[1].text = "• 营收增长 25%\n• 用户增长 40%\n• 市场份额提升 15%"
        
        # 结束页
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        textbox = slide.shapes.add_textbox(Inches(4), Inches(3), Inches(5), Inches(1))
        tf = textbox.text_frame
        tf.text = "感谢观看"
        p = tf.paragraphs[0]
        p.font.size = Pt(44)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        business_path = self.test_files_dir / "mock_business.pptx"
        prs.save(business_path)
        self.mock_files["business"] = business_path
        
        # 创建带图表的模板
        prs2 = Presentation()
        slide = prs2.slides.add_slide(prs2.slide_layouts[5])  # 空白布局
        
        # 添加表格
        table = slide.shapes.add_table(3, 3, Inches(1), Inches(1), Inches(8), Inches(4)).table
        table.cell(0, 0).text = "指标"
        table.cell(0, 1).text = "Q1"
        table.cell(0, 2).text = "Q2"
        table.cell(1, 0).text = "营收"
        table.cell(1, 1).text = "100万"
        table.cell(1, 2).text = "150万"
        
        # 添加形状
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(5.5), Inches(3), Inches(1))
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(0x00, 0x7A, 0xCC)
        
        chart_path = self.test_files_dir / "mock_chart.pptx"
        prs2.save(chart_path)
        self.mock_files["chart"] = chart_path
        
        print(f"  ✅ 创建模拟文件: {business_path.name}")
        print(f"  ✅ 创建模拟文件: {chart_path.name}")
    
    def get_mock_file(self, file_type: str = "business") -> Path:
        return self.mock_files.get(file_type, self.mock_files["business"])
    
    def setup_patches(self):
        """设置模拟补丁"""
        patches = []
        
        # 模拟鉴权
        auth_patcher = patch('feishu.auth.FeishuAuth._fetch_token')
        mock_auth = auth_patcher.start()
        mock_auth.return_value = "mock_token_12345"
        patches.append(auth_patcher)
        
        # 模拟get_token方法
        get_token_patcher = patch('feishu.auth.FeishuAuth.get_token')
        mock_get_token = get_token_patcher.start()
        mock_get_token.return_value = "mock_token_12345"
        patches.append(get_token_patcher)
        
        # 模拟获取文件信息
        info_patcher = patch('feishu.file_manager.FeishuFileManager.get_ppt_info')
        mock_info = info_patcher.start()
        mock_info.return_value = Mock(
            file_token="mock_file_token_123",
            file_name="mock_template.pptx",
            file_type="pptx",
            owner="测试用户",
            create_time="2026-04-03T10:00:00",
            version="1",
            size=102400,
            url="https://test.feishu.cn/slides/sldcnTest123"
        )
        patches.append(info_patcher)
        
        # 模拟下载
        download_patcher = patch('feishu.file_manager.FeishuFileManager.download_ppt')
        mock_download = download_patcher.start()
        
        def mock_download_impl(url, **kwargs):
            from feishu.file_manager import DownloadResult
            mock_file = self.get_mock_file("business")
            return DownloadResult(
                file_path=str(mock_file),
                file_name="mock_template.pptx",
                file_size=mock_file.stat().st_size,
                mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        
        mock_download.side_effect = mock_download_impl
        patches.append(download_patcher)
        
        # 模拟上传
        upload_patcher = patch('feishu.file_manager.FeishuFileManager.upload_ppt')
        mock_upload = upload_patcher.start()
        
        def mock_upload_impl(path, **kwargs):
            from feishu.file_manager import UploadResult
            return UploadResult(
                file_token="mock_upload_token_456",
                file_name="rebuilt_ppt.pptx",
                file_type="pptx",
                feishu_url="https://test.feishu.cn/slides/sldcnRebuilt456"
            )
        
        mock_upload.side_effect = mock_upload_impl
        patches.append(upload_patcher)
        
        return patches


# ==================== 全链路测试类 ====================

class FullIntegrationTest:
    """全链路集成测试"""
    
    def __init__(self, config: TestConfig):
        self.config = config
        self.logger = TestLogger(verbose=config.verbose)
        self.report = TestReport(start_time=datetime.now())
        self.mock_env: Optional[MockFeishuEnvironment] = None
        self.patches: List = []
        
        # 确保输出目录存在
        self.output_dir = Path(config.output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
    
    def setup(self):
        """测试前准备"""
        self.logger.section("测试环境准备")
        
        if self.config.mode in ["mock", "hybrid"]:
            # 创建模拟环境
            test_files_dir = Path("./test_files")
            test_files_dir.mkdir(exist_ok=True)
            self.mock_env = MockFeishuEnvironment(test_files_dir)
            self.patches = self.mock_env.setup_patches()
            
            # 设置模拟环境变量
            os.environ["FEISHU_APP_ID"] = "mock_cli_xxxxxxxx"
            os.environ["FEISHU_APP_SECRET"] = "mock_secret_xxxxxxxx"
            
            self.logger.success("模拟环境初始化完成")
        
        if self.config.mode == "real":
            # 检查真实环境变量
            app_id = os.getenv("FEISHU_APP_ID")
            app_secret = os.getenv("FEISHU_APP_SECRET")
            
            if not app_id or not app_secret:
                raise ValueError(
                    "真实模式需要配置飞书凭证:\n"
                    "  set FEISHU_APP_ID=cli_xxx\n"
                    "  set FEISHU_APP_SECRET=xxx"
                )
            
            self.logger.success(f"飞书凭证已配置 (App ID: {app_id[:10]}...)")
    
    def teardown(self):
        """测试后清理"""
        # 停止所有补丁
        for patcher in self.patches:
            patcher.stop()
        
        # 清理模拟环境变量
        if self.config.mode in ["mock", "hybrid"]:
            os.environ.pop("FEISHU_APP_ID", None)
            os.environ.pop("FEISHU_APP_SECRET", None)
        
        self.logger.info("测试环境清理完成")
    
    def run_step(self, name: str, func: Callable, *args, **kwargs) -> TestStep:
        """执行单个测试步骤"""
        start = time.time()
        
        try:
            result = func(*args, **kwargs)
            duration = time.time() - start
            
            step = TestStep(
                name=name,
                status="success",
                duration=duration,
                message="执行成功",
                details=result if isinstance(result, dict) else {}
            )
            self.logger.success(f"{name} - 成功 ({duration:.2f}s)")
            return step
            
        except Exception as e:
            duration = time.time() - start
            error_msg = str(e)
            suggestion = self._get_error_suggestion(e)
            
            step = TestStep(
                name=name,
                status="failed",
                duration=duration,
                message="执行失败",
                error=error_msg,
                suggestion=suggestion
            )
            self.logger.error(f"{name} - 失败: {error_msg}")
            if suggestion:
                self.logger.info(f"建议: {suggestion}")
            return step
    
    def _get_error_suggestion(self, error: Exception) -> Optional[str]:
        """根据错误类型获取建议"""
        error_msg = str(error).lower()
        
        if "权限" in error_msg or "permission" in error_msg:
            return "请检查飞书应用是否已开通所需权限: drive:file:download, drive:file:upload"
        elif "不存在" in error_msg or "not found" in error_msg:
            return "请检查PPT链接是否有效，文件是否已被删除"
        elif "鉴权" in error_msg or "auth" in error_msg:
            return "请检查 FEISHU_APP_ID 和 FEISHU_APP_SECRET 是否正确"
        elif "网络" in error_msg or "connection" in error_msg:
            return "请检查网络连接，确认可以访问飞书API"
        elif "文件格式" in error_msg or "format" in error_msg:
            return "请确保上传的是有效的 .pptx 文件"
        elif "大小" in error_msg or "size" in error_msg:
            return "文件过大，请使用更小的PPT文件"
        else:
            return "请查看详细错误日志，或联系管理员"
    
    # ==================== 测试步骤 ====================
    
    def step_1_environment_check(self) -> TestStep:
        """步骤1: 环境检查"""
        def check():
            results = {
                "mode": self.config.mode,
                "python_version": f"{sys.version_info.major}.{sys.version_info.minor}.{sys.version_info.micro}",
                "output_dir": str(self.output_dir),
                "dependencies": {}
            }
            
            # 检查关键依赖
            deps = ["pptx", "requests"]
            for dep in deps:
                try:
                    __import__(dep)
                    results["dependencies"][dep] = "✓"
                except ImportError:
                    results["dependencies"][dep] = "✗"
            
            # 检查环境变量
            if self.config.mode == "real":
                results["feishu_app_id"] = os.getenv("FEISHU_APP_ID", "")[:10] + "..." if os.getenv("FEISHU_APP_ID") else "未设置"
                results["feishu_app_secret"] = "已设置" if os.getenv("FEISHU_APP_SECRET") else "未设置"
            
            return results
        
        return self.run_step("环境检查", check)
    
    def step_2_auth(self) -> TestStep:
        """步骤2: 鉴权"""
        def auth():
            if self.config.mode == "mock":
                # 模拟鉴权成功
                return {"auth_type": "mock", "token": "mock_token_12345"}
            
            # 真实鉴权
            from feishu.auth import init_auth, get_auth
            
            app_id = os.getenv("FEISHU_APP_ID")
            app_secret = os.getenv("FEISHU_APP_SECRET")
            
            init_auth(app_id=app_id, app_secret=app_secret)
            auth_instance = get_auth()
            token = auth_instance.get_token()
            
            return {
                "auth_type": "real",
                "token_prefix": token[:20] + "..." if token else None
            }
        
        return self.run_step("飞书鉴权", auth)
    
    def step_3_parse_url(self) -> TestStep:
        """步骤3: 解析PPT链接"""
        def parse():
            urls = extract_ppt_urls(self.config.test_ppt_url)
            
            if not urls:
                raise ValueError(f"无法从输入中提取PPT链接: {self.config.test_ppt_url}")
            
            return {
                "input": self.config.test_ppt_url,
                "extracted_urls": urls,
                "url_count": len(urls)
            }
        
        return self.run_step("解析PPT链接", parse)
    
    def step_4_get_file_info(self) -> TestStep:
        """步骤4: 获取文件信息"""
        def get_info():
            from feishu.file_manager import FeishuFileManager
            
            fm = FeishuFileManager()
            info = fm.get_ppt_info(self.config.test_ppt_url)
            
            return {
                "file_token": info.file_token,
                "file_name": info.file_name,
                "file_type": info.file_type,
                "owner": info.owner,
                "size": info.size,
                "size_formatted": self._format_file_size(info.size)
            }
        
        return self.run_step("获取文件信息", get_info)
    
    def step_5_download_ppt(self) -> TestStep:
        """步骤5: 下载PPT文件"""
        def download():
            # 混合模式：使用本地文件
            if self.config.mode == "hybrid" and self.config.local_input_path:
                input_path = Path(self.config.local_input_path)
                if not input_path.exists():
                    raise FileNotFoundError(f"本地文件不存在: {input_path}")
                
                return {
                    "file_path": str(input_path),
                    "file_name": input_path.name,
                    "file_size": input_path.stat().st_size,
                    "source": "local"
                }
            
            # 模拟或真实下载
            from feishu.file_manager import FeishuFileManager
            
            fm = FeishuFileManager()
            result = fm.download_ppt(self.config.test_ppt_url)
            
            return {
                "file_path": result.file_path,
                "file_name": result.file_name,
                "file_size": result.file_size,
                "source": "download"
            }
        
        return self.run_step("下载PPT文件", download)
    
    def step_6_analyze_template(self, download_result: Dict) -> TestStep:
        """步骤6: 解析模板"""
        def analyze_ppt():
            file_path = download_result.get("file_path")
            
            profile = analyze(file_path)
            
            return {
                "width_cm": round(profile.geometry.width_cm, 2),
                "height_cm": round(profile.geometry.height_cm, 2),
                "theme_colors_count": len(profile.theme_colors),
                "font_scheme": {
                    "major_latin": profile.font_scheme.major_latin,
                    "minor_latin": profile.font_scheme.minor_latin
                },
                "layouts_count": len(profile.all_layouts),
                "masters_count": len(profile.masters)
            }
        
        return self.run_step("解析模板特征", analyze_ppt)
    
    def step_7_build_ppt(self, analyze_result: Dict) -> TestStep:
        """步骤7: 重建PPT"""
        def build():
            # 重新解析以获取完整profile
            file_path = self.report.steps[4].details.get("file_path")  # step_5的结果
            profile = analyze(file_path)
            
            builder = PPTBuilder(profile, output_dir=str(self.output_dir))
            
            build_options = BuildOptions(
                slide_count=self.config.slide_count,
                primary_color=self.config.primary_color
            )
            
            result = builder.build(build_options)
            
            # 验证输出文件
            output_path = Path(result.output_path)
            if not output_path.exists():
                raise FileNotFoundError(f"生成的PPT文件不存在: {output_path}")
            
            return {
                "output_path": str(result.output_path),
                "slide_count": result.slide_count,
                "layout_usage": len(result.layout_usage),
                "warnings_count": len(result.warnings),
                "file_size": output_path.stat().st_size,
                "file_size_formatted": self._format_file_size(output_path.stat().st_size)
            }
        
        return self.run_step("重建PPT", build)
    
    def step_8_upload_ppt(self, build_result: Dict) -> TestStep:
        """步骤8: 上传PPT到飞书"""
        def upload():
            output_path = build_result.get("output_path")
            
            from feishu.file_manager import FeishuFileManager
            
            fm = FeishuFileManager()
            result = fm.upload_ppt(output_path)
            
            return {
                "file_token": result.file_token,
                "file_name": result.file_name,
                "feishu_url": result.feishu_url
            }
        
        return self.run_step("上传到飞书", upload)
    
    def step_9_verify_result(self, upload_result: Dict) -> TestStep:
        """步骤9: 验证结果"""
        def verify():
            # 检查输出链接
            feishu_url = upload_result.get("feishu_url")
            
            if not feishu_url:
                raise ValueError("上传结果中没有飞书链接")
            
            # 验证链接格式
            if "feishu.cn" not in feishu_url and "mock" not in feishu_url:
                raise ValueError(f"飞书链接格式不正确: {feishu_url}")
            
            return {
                "feishu_url": feishu_url,
                "url_valid": True,
                "can_access": self.config.mode == "mock"  # 模拟环境可以直接"访问"
            }
        
        return self.run_step("验证结果", verify)
    
    def _format_file_size(self, size_bytes: int) -> str:
        """格式化文件大小"""
        if size_bytes < 1024:
            return f"{size_bytes} B"
        elif size_bytes < 1024 * 1024:
            return f"{size_bytes / 1024:.1f} KB"
        else:
            return f"{size_bytes / 1024 / 1024:.1f} MB"
    
    # ==================== 主测试流程 ====================
    
    def run(self) -> TestReport:
        """运行完整测试流程"""
        try:
            self.setup()
            
            self.logger.section("开始全链路集成测试")
            self.logger.info(f"测试模式: {self.config.mode}")
            self.logger.info(f"输出目录: {self.config.output_dir}")
            
            # 步骤1: 环境检查
            self.logger.step(1, 9, "环境检查")
            step1 = self.step_1_environment_check()
            self.report.add_step(step1)
            
            # 步骤2: 鉴权
            self.logger.step(2, 9, "飞书鉴权")
            step2 = self.step_2_auth()
            self.report.add_step(step2)
            
            if step2.status == "failed":
                self.logger.warning("鉴权失败，跳过后续步骤")
                return self._finalize_report()
            
            # 步骤3: 解析URL
            self.logger.step(3, 9, "解析PPT链接")
            step3 = self.step_3_parse_url()
            self.report.add_step(step3)
            
            # 步骤4: 获取文件信息
            self.logger.step(4, 9, "获取文件信息")
            step4 = self.step_4_get_file_info()
            self.report.add_step(step4)
            
            # 步骤5: 下载PPT
            self.logger.step(5, 9, "下载PPT文件")
            step5 = self.step_5_download_ppt()
            self.report.add_step(step5)
            
            if step5.status == "failed":
                self.logger.warning("下载失败，跳过后续步骤")
                return self._finalize_report()
            
            # 步骤6: 解析模板
            self.logger.step(6, 9, "解析模板特征")
            step6 = self.step_6_analyze_template(step5.details)
            self.report.add_step(step6)
            
            # 步骤7: 重建PPT
            self.logger.step(7, 9, "重建PPT")
            step7 = self.step_7_build_ppt(step6.details)
            self.report.add_step(step7)
            
            if step7.status == "failed":
                self.logger.warning("重建失败，跳过后续步骤")
                return self._finalize_report()
            
            # 步骤8: 上传PPT
            self.logger.step(8, 9, "上传到飞书")
            step8 = self.step_8_upload_ppt(step7.details)
            self.report.add_step(step8)
            
            # 步骤9: 验证结果
            self.logger.step(9, 9, "验证结果")
            step9 = self.step_9_verify_result(step8.details)
            self.report.add_step(step9)
            
        except Exception as e:
            self.logger.error(f"测试流程异常: {e}")
            traceback.print_exc()
        finally:
            self.teardown()
        
        return self._finalize_report()
    
    def _finalize_report(self) -> TestReport:
        """完成测试报告"""
        self.report.finalize()
        return self.report
    
    def print_report(self):
        """打印测试报告"""
        self.logger.section("全链路测试报告")
        
        # 汇总信息
        summary = self.report.summary
        print(f"\n📊 测试汇总")
        print(f"  总步骤数: {summary['total_steps']}")
        print(f"  成功: {summary['success']} | 失败: {summary['failed']} | 跳过: {summary['skipped']}")
        print(f"  成功率: {summary['success_rate']}")
        print(f"  总耗时: {summary['total_duration']}")
        
        # 详细步骤
        print(f"\n📋 详细步骤")
        for i, step in enumerate(self.report.steps, 1):
            status_icon = {"success": "✅", "failed": "❌", "skipped": "⏭️"}.get(step.status, "❓")
            print(f"\n  {i}. {status_icon} {step.name}")
            print(f"     状态: {step.status} | 耗时: {step.duration:.2f}s")
            
            if step.details:
                print(f"     详情:")
                for key, value in step.details.items():
                    if isinstance(value, dict):
                        print(f"       {key}:")
                        for k, v in value.items():
                            print(f"         - {k}: {v}")
                    else:
                        print(f"       - {key}: {value}")
            
            if step.error:
                print(f"     错误: {step.error}")
            
            if step.suggestion:
                print(f"     建议: {step.suggestion}")
        
        # 优化建议
        print(f"\n💡 优化建议")
        suggestions = self._generate_suggestions()
        for suggestion in suggestions:
            print(f"  • {suggestion}")
        
        # 保存报告
        report_path = self.output_dir / f"integration_test_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json"
        with open(report_path, 'w', encoding='utf-8') as f:
            json.dump(self.report.to_dict(), f, ensure_ascii=False, indent=2)
        print(f"\n📝 详细报告已保存: {report_path}")
    
    def _generate_suggestions(self) -> List[str]:
        """生成优化建议"""
        suggestions = []
        
        # 根据失败步骤生成建议
        failed_steps = [s for s in self.report.steps if s.status == "failed"]
        
        if not failed_steps:
            suggestions.append("所有测试步骤均通过，系统运行正常！")
            
            # 性能建议
            total_time = self.report.total_duration
            if total_time > 30:
                suggestions.append(f"总耗时较长({total_time:.1f}s)，建议优化文件下载/上传速度")
            
            # 检查是否有警告
            for step in self.report.steps:
                if step.name == "重建PPT" and step.details.get("warnings_count", 0) > 0:
                    suggestions.append(f"PPT重建过程中产生 {step.details['warnings_count']} 个警告，建议检查模板兼容性")
        else:
            for step in failed_steps:
                if step.suggestion:
                    suggestions.append(f"[{step.name}] {step.suggestion}")
        
        # 通用建议
        if self.config.mode == "mock":
            suggestions.append("当前使用模拟模式测试，建议在真实环境中再次验证")
        
        return suggestions


# ==================== 命令行入口 ====================

def main():
    """主入口"""
    parser = argparse.ArgumentParser(
        description="PPT复刻助手 - 全链路集成测试",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
测试模式说明:
  mock    - 完全模拟，无需飞书环境，适合本地开发测试
  real    - 真实环境，需要配置飞书凭证，测试完整链路
  hybrid  - 混合模式，使用本地文件，跳过下载步骤

示例:
  # 模拟测试
  python tests/test_full_integration.py --mode mock
  
  # 真实环境测试
  set FEISHU_APP_ID=cli_xxx
  set FEISHU_APP_SECRET=xxx
  python tests/test_full_integration.py --mode real --url "https://xxx.feishu.cn/slides/xxx"
  
  # 混合模式（使用本地文件）
  python tests/test_full_integration.py --mode hybrid --input ./template.pptx
        """
    )
    
    parser.add_argument(
        "--mode", "-m",
        choices=["mock", "real", "hybrid"],
        default="mock",
        help="测试模式 (默认: mock)"
    )
    
    parser.add_argument(
        "--url", "-u",
        default="https://test.feishu.cn/slides/sldcnTest123",
        help="飞书PPT链接 (real模式需要有效链接)"
    )
    
    parser.add_argument(
        "--input", "-i",
        help="本地PPT文件路径 (hybrid模式使用)"
    )
    
    parser.add_argument(
        "--output", "-o",
        default="./output/test_integration",
        help="输出目录 (默认: ./output/test_integration)"
    )
    
    parser.add_argument(
        "--pages", "-p",
        type=int,
        default=5,
        help="生成PPT的页面数 (默认: 5)"
    )
    
    parser.add_argument(
        "--color", "-c",
        help="指定主色调 (十六进制，如 FF5500)"
    )
    
    parser.add_argument(
        "--quiet", "-q",
        action="store_true",
        help="静默模式，减少输出"
    )
    
    args = parser.parse_args()
    
    # 创建测试配置
    config = TestConfig(
        mode=args.mode,
        test_ppt_url=args.url,
        local_input_path=args.input,
        output_dir=args.output,
        slide_count=args.pages,
        primary_color=args.color,
        verbose=not args.quiet
    )
    
    # 运行测试
    print("\n" + "=" * 60)
    print("  PPT复刻助手 - 全链路集成测试")
    print("=" * 60)
    
    try:
        test = FullIntegrationTest(config)
        report = test.run()
        test.print_report()
        
        # 根据结果返回退出码
        success_count = report.summary.get("success", 0)
        total_count = report.summary.get("total_steps", 0)
        
        if success_count == total_count:
            print("\n🎉 所有测试通过！\n")
            return 0
        else:
            print(f"\n⚠️  {total_count - success_count} 个步骤失败\n")
            return 1
            
    except Exception as e:
        print(f"\n💥 测试执行失败: {e}")
        traceback.print_exc()
        return 2


if __name__ == "__main__":
    sys.exit(main())
